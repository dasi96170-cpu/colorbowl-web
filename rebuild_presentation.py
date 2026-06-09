import os
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE

# 配色定義
NAVY = RGBColor(0x13, 0x27, 0x48)      # 墨綠深藍色 (主文字與主線條)
YELLOW = RGBColor(0xFD, 0xD1, 0x0E)    # 鵝黃色 (封面底色、強調背景)
CORAL = RGBColor(0xF2, 0x5B, 0x3B)     # 珊瑚橘紅 (主強調色、頂部粗飾線)
PAPER = RGBColor(0xF9, 0xF1, 0xEA)     # 米白色紙張 (投影片大背景)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # 白色 (卡片背景)
LINE = RGBColor(0xE7, 0xDE, 0xD2)      # 淺灰色細邊線 (卡片外框)
MUTED = RGBColor(0x6D, 0x75, 0x6F)     # 灰色 (副標與小字)
GREEN = RGBColor(0x2F, 0x8D, 0x46)     # 酪梨綠 (輔助色)

# 檔案路徑設定
base_dir = r"c:\Users\賴駿富\Desktop\AI數位內容與智慧流程設計實務班\期末專題"
source_path = os.path.join(base_dir, "彩碗_30_Demo簡報_v2_4.pptx")
dest_path = os.path.join(base_dir, r"網頁優化後\02_新版\彩碗_30_Demo簡報_v2_5_品牌重塑版.pptx")
new_logo_path = os.path.join(base_dir, r"logo\彩碗新logo.png")
bowl_mark_path = os.path.join(base_dir, r"outputs\colorbowl_brand_refresh\assets\colorbowl_bowl_mark_clear.png")

print(f"Source: {source_path}")
print(f"Destination: {dest_path}")

# 確保輸出目錄存在
dest_dir = os.path.dirname(dest_path)
if not os.path.exists(dest_dir):
    os.makedirs(dest_dir)

# ----------------------------------------------------------------------
# 輔助排版與繪製函數
# ----------------------------------------------------------------------

# 繪製矩形
def DrawRect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    shape.name = "RebuiltShape"
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

# 繪製圓角卡片
def DrawCard(slide, x, y, w, h, fill_color, line_color=None, top_bar_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
    card.name = "RebuiltCard"
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if line_color is not None:
        card.line.color.rgb = line_color
        card.line.width = Pt(1.0)
    else:
        card.line.fill.background()
    
    # 在卡片頂部加裝飾橘色線條 (向左右內縮 2 點以契合圓角)
    if top_bar_color is not None:
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x + 2), Pt(y), Pt(w - 4), Pt(5))
        top_bar.name = "RebuiltShape"
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = top_bar_color
        top_bar.line.fill.background()
    return card

# 繪製橫線/直線/斜線
def DrawLine(slide, x1, y1, x2, y2, color, weight=1.5):
    # 如果是水平線
    if abs(y1 - y2) < 0.1:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x1), Pt(y1), Pt(abs(x2 - x1)), Pt(weight))
        shape.name = "RebuiltShape"
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    # 如果是垂直線
    elif abs(x1 - x2) < 0.1:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x1), Pt(y1), Pt(weight), Pt(abs(y2 - y1)))
        shape.name = "RebuiltShape"
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape
    # 斜線
    else:
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
        connector.name = "RebuiltShape"
        connector.line.color.rgb = color
        connector.line.width = Pt(weight)
        return connector

# 繪製圖片
def DrawPicture(slide, path, x, y, w, h):
    if os.path.exists(path):
        try:
            pic = slide.shapes.add_picture(path, Pt(x), Pt(y), Pt(w), Pt(h))
            pic.name = "RebuiltPicture"
            return pic
        except Exception as e:
            print(f"Error adding picture {path}: {e}")
    return None

# 檢查 Shape 是否有文字
def HasText(shape):
    if not shape.has_text_frame:
        return False
    if not shape.text_frame.text.strip():
        return False
    return True

# 將文字 Shape 移到最上層
def bring_to_front(shape):
    try:
        spTree = shape._element.getparent()
        spTree.remove(shape._element)
        spTree.append(shape._element)
    except Exception:
        pass

def PushTextToFront(slide):
    text_shapes = []
    for shape in list(slide.shapes):
        if HasText(shape):
            text_shapes.append(shape)
    for shape in text_shapes:
        bring_to_front(shape)

# 套用通用的頁面背景與標頭樣式
def ApplyCommonTheme(slide, PAPER, YELLOW, NAVY, CORAL, bowl_mark_path):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    
    # 頂部導覽列黃色長條
    DrawRect(slide, 0, 0, 960, 106, YELLOW)
    
    # 左側邊緣細橘條
    DrawRect(slide, 0, 0, 11, 540, CORAL)
    
    # 左上角墨綠色小方塊 (LOGO後面)
    DrawRect(slide, 38, 30, 58, 54, NAVY)
    
    # 橫貫頂部的墨綠色細裝飾線
    DrawLine(slide, 42, 105, 916, 105, NAVY, 1.6)
    
    # 左側橘色強效雙粗飾線
    DrawLine(slide, 42, 108, 210, 108, CORAL, 3.2)
    
    # 右上角去背 Logo 碗形圖示
    if os.path.exists(bowl_mark_path):
        DrawPicture(slide, bowl_mark_path, 856, 18, 70, 45)

# 統一設定文字格式與配色
def FormatText(shape, color, size, bold, fontName="Microsoft JhengHei UI"):
    try:
        # 移除背景與框線，只保留文字
        shape.fill.background()
        shape.line.fill.background()
        
        tf = shape.text_frame
        tf.margin_left = Pt(3)
        tf.margin_right = Pt(3)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        
        for paragraph in tf.paragraphs:
            paragraph.font.name = fontName
            paragraph.font.size = Pt(size)
            paragraph.font.color.rgb = color
            paragraph.font.bold = bold
            
            # 設定 runs 屬性，確保個別樣式也被覆蓋
            for run in paragraph.runs:
                run.font.name = fontName
                run.font.size = Pt(size)
                run.font.color.rgb = color
                run.font.bold = bold
    except Exception as e:
        pass

# 遍歷格式化每一頁的所有文字框
def StyleAllSlideText(slide, slide_index, NAVY, CORAL, MUTED, WHITE, YELLOW, GREEN):
    for shape in list(slide.shapes):
        if not HasText(shape):
            continue
        x = shape.left.pt
        y = shape.top.pt
        w = shape.width.pt
        txt = shape.text_frame.text.strip()
        
        # 預設通用文字格式
        FormatText(shape, NAVY, 12.5, False)
        
        # 依頁碼特化文字大小與配色
        if slide_index == 1: # 封面頁
            if y < 130:
                FormatText(shape, NAVY, 18, True)
            elif y < 220:
                FormatText(shape, NAVY, 34, True)
            elif y < 300:
                FormatText(shape, NAVY, 18, True)
            elif y > 320:
                FormatText(shape, WHITE, 11.5, False)
            continue
        
        if slide_index == 10: # 新增的 Logo 視覺解析頁
            if x < 90 and y < 95:
                FormatText(shape, WHITE, 16, True)
            elif x > 90 and y < 58:
                FormatText(shape, NAVY, 22, True)
            elif x > 90 and y < 100:
                FormatText(shape, CORAL, 10.5, True)
            elif y > 492:
                FormatText(shape, MUTED, 8.5, False)
            elif "01_" in txt or "02_" in txt:
                FormatText(shape, NAVY, 18, True)
            elif any(txt.startswith(prefix) for prefix in ["明亮黃", "珊瑚橘", "深藍色", "酪梨綠", "紫甘藍"]):
                FormatText(shape, CORAL, 13, True)
            else:
                FormatText(shape, NAVY, 12, False)
            continue
        
        if slide_index == 12: # 落地藍圖過渡頁 (原 Slide 11)
            if y < 120:
                FormatText(shape, CORAL, 18, True)
            elif y < 280:
                FormatText(shape, NAVY, 34, True)
            elif y > 470:
                FormatText(shape, MUTED, 10, False)
            else:
                FormatText(shape, NAVY, 14, False)
            continue
        
        if slide_index == 13: # 結尾頁 (原 Slide 12)
            if y < 230:
                FormatText(shape, YELLOW, 24, True)
            elif y < 340:
                FormatText(shape, WHITE, 34, True)
            elif y > 390:
                FormatText(shape, NAVY, 15, False)
            else:
                FormatText(shape, WHITE, 15, False)
            continue
        
        # 內容頁的一般文字區分規則
        if x < 90 and y < 95:
            FormatText(shape, WHITE, 16, True) # 頂部大數字標示 (如 01, 02)
        elif x > 90 and y < 58:
            FormatText(shape, NAVY, 22, True) # 頂部大標題
        elif x > 90 and y < 100:
            FormatText(shape, CORAL, 10.5, True) # 頂部英文小副標
        elif y > 492:
            FormatText(shape, MUTED, 8.5, False) # 底部註記來源
        elif re.match(r"^(PERSONA|STEP|BEFORE|AFTER|WHY|HOW|WHAT|Pokeworks|S 優勢|W 劣勢|O 機會|T 威脅)", txt):
            FormatText(shape, NAVY, 14.5, True) # 卡片內的小標題
        elif re.match(r"^(01|02|03|1|2|3|A[1-5])$", txt):
            FormatText(shape, WHITE, 17, True) # 圓形大數字符號
        elif txt.startswith("“") or txt.startswith("”") or txt.startswith(chr(0x201C)):
            FormatText(shape, CORAL, 16, True) # 金句引用括號
        elif w > 750 and y < 170:
            FormatText(shape, NAVY, 13.5, False) # 段落引言
        elif w > 750:
            FormatText(shape, NAVY, 12.2, False) # 寬度大的文字
        else:
            FormatText(shape, NAVY, 12.5, False) # 標準內文

# 刪除 shape 輔助函數
def delete_shape(shape):
    try:
        sp = shape._element
        sp.getparent().remove(sp)
    except Exception as e:
        print(f"Error deleting shape: {e}")

# 移動投影片輔助函數
def move_slide(prs, slide, target_idx):
    sldIdLst = prs.slides._sldIdLst
    slide_id = slide.slide_id
    slide_element = None
    for sldId in sldIdLst:
        if sldId.id == slide_id:
            slide_element = sldId
            break
    if slide_element is not None:
        sldIdLst.remove(slide_element)
        sldIdLst.insert(target_idx, slide_element)

# ----------------------------------------------------------------------
# 主執行邏輯
# ----------------------------------------------------------------------

# 讀取投影片
pres = Presentation(source_path)
print(f"Original Slide Count: {len(pres.slides)}")

# ----------------------------------------------------------------------
# 第一部分：在 Slide 9 之後插入全新的 Logo 視覺分析投影片 (成為 Slide 10)
# ----------------------------------------------------------------------

# 尋找空白 layout
blank_layout = None
for layout in pres.slide_layouts:
    if len(layout.placeholders) == 0:
        blank_layout = layout
        break
if not blank_layout:
    blank_layout = pres.slide_layouts[6]

slide10 = pres.slides.add_slide(blank_layout)
move_slide(pres, slide10, 9) # 移至第 10 頁 (0-based index 為 9)
print("Created New Slide 10 and moved to position 10")

# 填寫 Slide 10 內容 (使用 TextBox 來添加文字)
# 1. 頂部大數字 "A2"
tbNum = slide10.shapes.add_textbox(Pt(38), Pt(20), Pt(50), Pt(50))
tbNum.text_frame.text = "A2"

# 2. 標題與副標
tbTitle = slide10.shapes.add_textbox(Pt(100), Pt(20), Pt(600), Pt(80))
tbTitle.text_frame.text = "品牌分析｜新版 Logo 視覺解析\nVisual identity — repositioning the brand"

# 3. 核心分析
tbDesc = slide10.shapes.add_textbox(Pt(62), Pt(126), Pt(832), Pt(70))
tbDesc.text_frame.text = "新版 Logo 以鮮明黃色作為背景，搭配珊瑚橘碗身、鮭魚、酪梨、米飯與紫甘藍等食材元素，強化品牌「健康、繽紛、新鮮」的餐飲印象。整體視覺比舊版更直接傳達 Poke 碗餐特色，也更符合年輕族群與外帶健康餐市場的品牌語言。"

# 4. 主要顏色與色彩對標 (製作左右兩欄白色圓角卡片)
DrawCard(slide10, 62, 216, 400, 240, WHITE, LINE, CORAL)
tbColorTitle = slide10.shapes.add_textbox(Pt(80), Pt(226), Pt(360), Pt(30))
tbColorTitle.text_frame.text = "01_色彩計畫與視覺印象"

tbColors = slide10.shapes.add_textbox(Pt(80), Pt(266), Pt(360), Pt(180))
tbColors.text_frame.text = "明亮黃色背景 (#FDD10E) — 活潑年輕、高度吸睛 (83.8%)\n\n珊瑚橘紅碗身 (#F25B3B) — 食慾、熱情、品牌主體 (7.9%)\n\n深海軍藍文字 (#132748) — 穩重對比、提高辨識 (2.1%)\n\n米飯白與繽紛色 — 乾淨主食感與多樣食材新鮮感 (6.2%)"

# 5. 食材與符號解析
DrawCard(slide10, 494, 216, 400, 240, WHITE, LINE, CORAL)
tbFoodTitle = slide10.shapes.add_textbox(Pt(512), Pt(226), Pt(360), Pt(30))
tbFoodTitle.text_frame.text = "02_食材辨識與品牌定位"

tbFoods = slide10.shapes.add_textbox(Pt(512), Pt(266), Pt(360), Pt(180))
tbFoods.text_frame.text = "鮭魚與米飯 — 突出 Poke 核心靈魂，明確傳達產品定位。\n\n酪梨與紫甘藍 — 深綠、淺綠與紫紅的漸層配搭，展現「彩碗」的健康、新鮮與彩色概念。\n\n深藍筷子與字體 — 筷子帶來明確的餐飲感，寬字距海軍藍字體展現俐落的都會感。"

# 6. 底部來源註解
tbSource = slide10.shapes.add_textbox(Pt(62), Pt(500), Pt(832), Pt(30))
tbSource.text_frame.text = "Logo 色彩比例依據圖面取樣統計，詳見《品牌識別系統規劃書 v1.0》。"

# ----------------------------------------------------------------------
# 第二部分：遍歷並重繪所有 18 頁投影片的視覺元素與排版
# ----------------------------------------------------------------------

for idx, slide in enumerate(pres.slides):
    slide_index = idx + 1
    print(f"Processing Slide {slide_index}...")
    
    # 1. 刪除原有的裝飾形狀 (只保留有文字的框框與原有的圖示/表格/圖片，清除舊的底色塊以防衝突)
    shapes_to_delete = []
    for shape in list(slide.shapes):
        if not HasText(shape):
            st = int(shape.shape_type)
            # 我們要保留的類型 (13: PICTURE, 19: TABLE, 7: OLE_OBJECT, 3: CHART, 6: GROUP)
            keep_types = [13, 19, 7, 3, 6]
            if st not in keep_types:
                name = shape.name.lower()
                # 保留有特定命名的 Shape 或原有的食材Icon圖標
                if not any(x in name for x in ["picture", "table", "ole", "diagram"]):
                    shapes_to_delete.append(shape)
                    
    for shape in shapes_to_delete:
        delete_shape(shape)
        
    # 2. 針對不同的 Slide Index 進行特化排版 (繪製白色圓角卡片與裝飾)
    if slide_index == 1:
        # 封面頁：大面積金黃色背景
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = YELLOW
        
        # 墨綠深藍色大 Banner
        DrawRect(slide, 0, 424, 960, 116, NAVY)
        # 珊瑚橘細線
        DrawRect(slide, 0, 414, 960, 10, CORAL)
        
        # 畫一條珊瑚橘大橫線與綠色短線 (裝飾)
        DrawLine(slide, 68, 306, 318, 306, CORAL, 3)
        DrawLine(slide, 68, 313, 208, 313, GREEN, 2)
        
        # 右側置入新版大 Logo (logo\彩碗新logo.png)
        if os.path.exists(new_logo_path):
            DrawPicture(slide, new_logo_path, 600, 68, 280, 280)
            
    elif slide_index == 12:
        # 落地藍圖過渡頁 (原 Slide 11)：大米白色底 + 金黃底色塊 8% 透明
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PAPER
        DrawRect(slide, 0, 0, 960, 540, YELLOW) # 用金黃色底板
        # 左側細橘條
        DrawRect(slide, 0, 0, 16, 540, CORAL)
        # 圓角大卡片
        DrawCard(slide, 62, 286, 832, 150, WHITE, LINE, None)
        DrawLine(slide, 70, 276, 300, 276, CORAL, 3)
        # 右側置入去背 Logo
        if os.path.exists(bowl_mark_path):
            DrawPicture(slide, bowl_mark_path, 640, 80, 250, 160)
            
    elif slide_index == 13:
        # 結尾頁 (原 Slide 12)：深海軍藍大背景，極致奢華
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = NAVY
        DrawRect(slide, 0, 392, 960, 148, YELLOW)
        DrawRect(slide, 0, 0, 18, 540, CORAL)
        DrawLine(slide, 72, 302, 302, 302, CORAL, 3)
        # 中央偏右大去背 Logo
        if os.path.exists(bowl_mark_path):
            DrawPicture(slide, bowl_mark_path, 646, 90, 250, 160)
            
    else:
        # 通用內容頁 (套用通用背景樣式)
        ApplyCommonTheme(slide, PAPER, YELLOW, NAVY, CORAL, bowl_mark_path)
        
        # 依頁面索引特化內容卡片排版
        if slide_index == 2: # 提案總覽 (三列卡片)
            DrawCard(slide, 43, 187, 274, 252, WHITE, LINE, CORAL)
            DrawCard(slide, 338, 187, 274, 252, WHITE, LINE, CORAL)
            DrawCard(slide, 634, 187, 274, 252, WHITE, LINE, CORAL)
        elif slide_index == 3: # 市場分析背景 (垂直多列卡片)
            DrawCard(slide, 43, 161, 878, 92, WHITE, LINE, None)
            DrawCard(slide, 43, 266, 878, 92, WHITE, LINE, None)
            DrawCard(slide, 43, 371, 878, 92, WHITE, LINE, None)
            DrawCard(slide, 42, 492, 880, 34, WHITE, LINE, None)
        elif slide_index == 4: # 市場分析競品 (大卡片 + 頂部深藍底 + 底部黃色強調)
            DrawCard(slide, 43, 160, 878, 250, WHITE, LINE, None)
            DrawRect(slide, 50, 166, 865, 36, NAVY)
            DrawRect(slide, 50, 356, 865, 53, YELLOW, None)
            DrawCard(slide, 50, 424, 864, 68, WHITE, LINE, None)
        elif slide_index == 5: # 目標 Persona (三列卡片 + 底部黃色強調)
            DrawCard(slide, 43, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 338, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 634, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 43, 364, 871, 112, WHITE, LINE, None)
            DrawRect(slide, 43, 364, 871, 112, YELLOW, None)
        elif slide_index == 6: # SWOT 交叉策略 (四格卡片與對齊)
            DrawCard(slide, 43, 169, 166, 94, WHITE, LINE, None)
            DrawCard(slide, 223, 169, 691, 94, WHITE, LINE, None)
            DrawCard(slide, 43, 277, 166, 94, WHITE, LINE, None)
            DrawCard(slide, 223, 277, 691, 94, WHITE, LINE, None)
            DrawCard(slide, 43, 385, 166, 94, WHITE, LINE, None)
            DrawCard(slide, 223, 385, 691, 94, WHITE, LINE, None)
            DrawRect(slide, 43, 169, 166, 310, NAVY)
        elif slide_index == 7: # 策略驗證 (三列卡片，對標高亮)
            DrawCard(slide, 43, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 338, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 634, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 43, 364, 871, 112, WHITE, LINE, None)
            DrawRect(slide, 338, 130, 274, 216, YELLOW, None)
        elif slide_index == 8: # 黃金圈轉型 (三列卡片，高亮黃金圈WHY)
            DrawCard(slide, 43, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 338, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 634, 130, 274, 216, WHITE, LINE, CORAL)
            DrawCard(slide, 43, 364, 871, 112, WHITE, LINE, None)
            DrawRect(slide, 634, 130, 274, 216, YELLOW, None)
        elif slide_index == 9: # 展示段實作進度 (雙大卡)
            DrawCard(slide, 43, 173, 425, 281, WHITE, LINE, CORAL)
            DrawCard(slide, 490, 173, 425, 281, WHITE, LINE, CORAL)
            DrawRect(slide, 43, 173, 425, 281, YELLOW, None)
        elif slide_index == 11: # 落地行銷漏斗 (雙大卡 + 底部卡片)
            DrawCard(slide, 43, 162, 428, 284, WHITE, LINE, CORAL)
            DrawCard(slide, 486, 162, 428, 284, WHITE, LINE, CORAL)
            DrawCard(slide, 43, 461, 871, 29, WHITE, LINE, None)
        elif slide_index == 14: # 視覺風格指南 (左四小卡 + 右側大卡)
            DrawCard(slide, 50, 126, 421, 133, WHITE, LINE, None)
            DrawCard(slide, 486, 126, 421, 133, WHITE, LINE, None)
            DrawCard(slide, 50, 270, 421, 133, WHITE, LINE, None)
            DrawCard(slide, 486, 270, 421, 133, WHITE, LINE, None)
            DrawCard(slide, 50, 428, 864, 61, WHITE, LINE, None)
            # 四小卡左邊緣各加上對應顏色飾條
            DrawRect(slide, 50, 126, 421, 8, GREEN)
            DrawRect(slide, 486, 126, 421, 8, CORAL)
            DrawRect(slide, 50, 270, 421, 8, YELLOW)
            DrawRect(slide, 486, 270, 421, 8, NAVY)
        elif slide_index == 15: # 資料來源 (左四小卡 + 右四小卡)
            DrawCard(slide, 43, 173, 421, 122, WHITE, LINE, None)
            DrawCard(slide, 482, 173, 421, 122, WHITE, LINE, None)
            DrawCard(slide, 43, 313, 421, 122, WHITE, LINE, None)
            DrawCard(slide, 482, 313, 421, 122, WHITE, LINE, None)
            # 四小卡左邊緣各加上對應顏色飾條
            DrawRect(slide, 43, 173, 9, 122, GREEN)
            DrawRect(slide, 482, 173, 9, 122, CORAL)
            DrawRect(slide, 43, 313, 9, 122, YELLOW)
            DrawRect(slide, 482, 313, 9, 122, NAVY)
        elif slide_index == 16: # 結語前頁面 (六小卡並排)
            DrawCard(slide, 43, 166, 274, 130, WHITE, LINE, CORAL)
            DrawCard(slide, 338, 166, 274, 130, WHITE, LINE, CORAL)
            DrawCard(slide, 634, 166, 274, 130, WHITE, LINE, CORAL)
            DrawCard(slide, 43, 313, 274, 130, WHITE, LINE, CORAL)
            DrawCard(slide, 338, 313, 274, 130, WHITE, LINE, CORAL)
            DrawCard(slide, 634, 313, 274, 130, WHITE, LINE, CORAL)
        elif slide_index == 17: # 網站 Before & After (三大橫向重疊卡片)
            DrawCard(slide, 43, 166, 878, 306, WHITE, LINE, None)
            DrawRect(slide, 43, 166, 936, 90, CORAL, None)
            DrawRect(slide, 137, 274, 749, 90, YELLOW, None)
            DrawRect(slide, 230, 382, 562, 90, GREEN, None)
            DrawLine(slide, 102, 256, 756, 382, CORAL, 2.4)
        elif slide_index == 18: # LINE Bot AI配餐 (三行大卡)
            DrawCard(slide, 43, 166, 274, 259, WHITE, LINE, CORAL)
            DrawCard(slide, 338, 166, 274, 259, WHITE, LINE, CORAL)
            DrawCard(slide, 634, 166, 274, 259, WHITE, LINE, CORAL)
            DrawRect(slide, 43, 166, 274, 50, YELLOW, None)
            DrawRect(slide, 338, 166, 274, 50, CORAL, None)
            DrawRect(slide, 634, 166, 274, 50, GREEN, None)

    # 3. 對該頁的所有文字進行格式套用
    StyleAllSlideText(slide, slide_index, NAVY, CORAL, MUTED, WHITE, YELLOW, GREEN)
    
    # 4. 微調第 6 頁 (SWOT交叉) 左側深藍卡片上的文字顏色為白色，以利閱讀
    if slide_index == 6:
        for sh in list(slide.shapes):
            if HasText(sh) and sh.left.pt < 210 and sh.top.pt > 160:
                FormatText(sh, WHITE, 15, True)
                
    # 5. 強制將文字 Shape 堆疊到最上層，防止被我們剛剛繪製的圓角卡片遮擋
    PushTextToFront(slide)

# 保存重建後的簡報
pres.save(dest_path)
print(f"Saved Rebuilt Presentation to: {dest_path}")
print("Process completed successfully!")
